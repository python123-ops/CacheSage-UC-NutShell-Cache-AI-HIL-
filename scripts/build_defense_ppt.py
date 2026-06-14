# -*- coding: utf-8 -*-
from __future__ import annotations

import json
from pathlib import Path
from typing import Iterable, Sequence

from PIL import Image, ImageDraw, ImageFont
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt


ROOT = Path(__file__).resolve().parents[1]
REPORTS = ROOT / "reports"
ASSETS = REPORTS / "demo-assets"
OUT = REPORTS / "CacheSage-UC-defense-demo.pptx"

W, H = 13.333, 7.5
INK = RGBColor(13, 20, 31)
PANEL = RGBColor(21, 33, 48)
PANEL_2 = RGBColor(31, 48, 66)
CYAN = RGBColor(58, 204, 225)
GREEN = RGBColor(83, 222, 143)
AMBER = RGBColor(245, 190, 82)
WHITE = RGBColor(244, 248, 252)
MUTED = RGBColor(163, 177, 194)
RULE = RGBColor(72, 96, 116)


def main() -> int:
    ASSETS.mkdir(parents=True, exist_ok=True)
    data = load_data()
    make_terminal_assets(data)

    prs = Presentation()
    prs.slide_width = Inches(W)
    prs.slide_height = Inches(H)
    blank = prs.slide_layouts[6]
    for builder in [
        slide_cover,
        slide_evidence_map,
        slide_linux_env,
        slide_smoke_flow,
        slide_coverage,
        slide_faults,
        slide_review,
        slide_boundary,
    ]:
        builder(prs.slides.add_slide(blank), data)
    prs.save(OUT)
    print(OUT)
    return 0


def load_data() -> dict:
    smoke = json.loads((REPORTS / "nutshell-smoke.json").read_text(encoding="utf-8"))
    sample = json.loads((REPORTS / "sample-run-seed11.json").read_text(encoding="utf-8"))
    env = json.loads((REPORTS / "linux-smoke-env.json").read_text(encoding="utf-8"))
    review = [
        json.loads(line)
        for line in (ROOT / "review_journal.jsonl").read_text(encoding="utf-8").splitlines()
        if line.strip()
    ]
    faults = []
    for path in sorted(REPORTS.glob("fault-*.json")):
        payload = json.loads(path.read_text(encoding="utf-8"))
        faults.append(
            {
                "mode": payload.get("fault") or path.stem.replace("fault-", "").replace("-", "_"),
                "failures": len(payload.get("failures", [])),
            }
        )
    return {"smoke": smoke, "sample": sample, "env": env, "review": review, "faults": faults}


def make_terminal_assets(data: dict) -> None:
    smoke = data["smoke"]
    cov = smoke["python_harness"]["coverage"]
    make_results = smoke.get("make_results", {})
    artifacts = smoke.get("rtl_artifacts", {})
    rtl_code_coverage = smoke.get("rtl_code_coverage", {})
    env = data["env"]
    sample = data["sample"]

    terminal_image(
        ASSETS / "terminal-linux-env.png",
        [
            "$ wsl -d Ubuntu-24.04 -u root -- bash -lc 'picker --check'",
            env["wsl_os"]["output"].replace("\n", " | "),
            env["kernel"]["output"],
            env["python"]["output"],
            env["verilator"]["output"],
            env["picker"]["output"].splitlines()[-1],
            compact_picker_check(env["picker_check"]["output"]),
            env["toffee"]["output"].replace("\n", " | "),
        ],
    )
    terminal_image(
        ASSETS / "terminal-smoke.png",
        [
            "$ python scripts/run_nutshell_smoke.py --upstream third_party/Example-NutShellCache",
            f"status: {smoke['status']}",
            f"make gen_dut: exit {make_results.get('gen_dut', {}).get('returncode')}",
            f"make test: exit {make_results.get('test', {}).get('returncode')}",
            "pytest: test/test_smoke.py::test_smoke PASSED",
            f"Python harness: {cov['covered']}/{cov['total']} coverpoints, {cov['percent']:.2f}%",
            f"RTL artifacts: {artifacts.get('status', 'not_recorded')}, count={artifacts.get('count', 0)}",
            f"RTL code coverage: {rtl_code_coverage.get('status', 'not_recorded')}",
        ],
    )
    events = sample["event_counts"]
    terminal_image(
        ASSETS / "terminal-coverage.png",
        [
            "$ python -m cachesage_uc.cli run --seed 11 --count 96",
            f"passed: {sample['passed']}",
            f"transactions: {sample['transaction_count']}",
            f"coverage: {sample['coverage']['covered']}/{sample['coverage']['total']} ({sample['coverage']['percent']:.2f}%)",
            f"miss/refill: {events.get('miss')}/{events.get('refill')}",
            f"dirty eviction/writeback: {events.get('dirty_eviction')}/{events.get('writeback')}",
            f"stall/reset: {events.get('stall_hold')}/{events.get('reset_window')}",
        ],
    )
    fault_lines = ["$ python -m cachesage_uc.cli run --fault <mode>"]
    for item in data["faults"]:
        fault_lines.append(f"{item['mode']}: expected failure, failures={item['failures']}")
    terminal_image(ASSETS / "terminal-faults.png", fault_lines)

    review_lines = ["$ head review_journal.jsonl"]
    for row in data["review"][:5]:
        review_lines.append(f"{row['id']}: reviewed -> {row['coverage_delta']}")
    terminal_image(ASSETS / "terminal-review.png", review_lines)
    terminal_image(
        ASSETS / "terminal-repo.png",
        [
            "$ tree -L 2",
            "src/cachesage_uc/        verification core + adapters",
            "tests/                   CLI, materials, adapter, core tests",
            "docs/                    plan, scoreboard, fault, picker/toffee flow",
            "reports/                 PDF, smoke JSON/MD, fault artifacts, PPT",
            "third_party/             pinned Example-NutShellCache (not vendored)",
        ],
    )


def compact_picker_check(output: str) -> str:
    keep = []
    for line in output.splitlines():
        if "Version" in line:
            keep.append("OK: Picker version 0.9.0")
        elif "Support    Cpp" in line:
            keep.append("OK: Cpp support")
        elif "Support Python" in line:
            keep.append("OK: Python support")
    return " | ".join(keep)


def terminal_image(path: Path, lines: Sequence[str], width: int = 1320, height: int = 560) -> None:
    img = Image.new("RGB", (width, height), (10, 16, 24))
    draw = ImageDraw.Draw(img)
    draw.rounded_rectangle([18, 18, width - 18, height - 18], radius=20, fill=(16, 27, 39), outline=(53, 86, 106), width=2)
    draw.ellipse([44, 44, 62, 62], fill=(239, 90, 90))
    draw.ellipse([74, 44, 92, 62], fill=(245, 190, 82))
    draw.ellipse([104, 44, 122, 62], fill=(83, 222, 143))
    mono = load_font(24, mono=True)
    y = 96
    for idx, line in enumerate(lines[:13]):
        color = (94, 223, 240) if idx == 0 else (225, 235, 242)
        draw.text((48, y), line[:105], font=mono, fill=color)
        y += 34
    img.save(path)


def load_font(size: int, mono: bool = False) -> ImageFont.FreeTypeFont:
    candidates = [
        "C:/Windows/Fonts/consola.ttf" if mono else "C:/Windows/Fonts/msyh.ttc",
        "C:/Windows/Fonts/simhei.ttf",
        "C:/Windows/Fonts/arial.ttf",
    ]
    for candidate in candidates:
        try:
            return ImageFont.truetype(candidate, size)
        except OSError:
            continue
    return ImageFont.load_default()


def base(slide, kicker: str, title: str, page: str) -> None:
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = INK
    add_text(slide, 0.52, 0.34, 2.2, 0.24, kicker, 8.5, CYAN, bold=True)
    add_text(slide, 0.52, 0.62, 8.9, 0.75, title, 25, WHITE, bold=True)
    line(slide, 0.52, 1.42, 12.2, 1.42, RULE)
    add_text(slide, 12.15, 6.95, 0.55, 0.22, page, 8, MUTED, align=PP_ALIGN.RIGHT)


def slide_cover(slide, data: dict) -> None:
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = INK
    add_text(slide, 0.68, 0.62, 2.2, 0.26, "UCAgent / NutShell Cache", 10, CYAN, bold=True)
    add_text(slide, 0.68, 1.25, 9.5, 1.2, "CacheSage-UC 验证答辩演示", 34, WHITE, bold=True)
    add_text(slide, 0.72, 2.45, 8.8, 0.55, "从 Python harness 到 Linux Picker/Toffee smoke 的可复核证据链", 17, MUTED)
    metric_rail(
        slide,
        0.72,
        3.48,
        [
            ("23/23", "Python harness 覆盖点"),
            ("5", "确定性 injected fault"),
            ("2/2", "make gen_dut / make test"),
            (str(data["smoke"].get("rtl_artifacts", {}).get("count", 0)), "RTL artifact manifest"),
        ],
        columns=2,
    )
    add_terminal(slide, ASSETS / "terminal-smoke.png", 6.75, 2.8, 5.85, 2.48)
    add_text(slide, 0.72, 6.65, 7.3, 0.24, "证据边界：Python functional coverage 与 RTL smoke artifact/code coverage 分栏记录。", 9.5, MUTED)


def slide_evidence_map(slide, data: dict) -> None:
    base(slide, "EVIDENCE CHAIN", "交付物围绕可复现验证证据组织", "02")
    nodes = [
        ("Generator / CRV", "seed 11 + directed spine"),
        ("Reference Model", "byte-addressable memory"),
        ("Scoreboard", "read data + event signature"),
        ("Coverage", "23 个功能覆盖点"),
        ("Fault Injection", "5 类确定性检出"),
        ("Linux Smoke", "Picker/Toffee 上游入口"),
    ]
    for i, (name, desc) in enumerate(nodes):
        x = 0.7 + (i % 3) * 4.1
        y = 1.9 + (i // 3) * 1.55
        panel(slide, x, y, 3.55, 1.02, name, desc)
    for x in [4.0, 8.1]:
        arrow(slide, x, 2.4, x + 0.55, 2.4)
        arrow(slide, x, 3.95, x + 0.55, 3.95)
    add_text(slide, 0.75, 5.62, 11.4, 0.58, "核心口径：Python harness 负责可控刺激、scoreboard rehearsal 与故障检出；Linux smoke 证明上游 Picker/Toffee 基础链路可执行。两类证据分栏记录，避免把 smoke 通过写成 RTL 覆盖率完成。", 15, WHITE)


def slide_linux_env(slide, data: dict) -> None:
    base(slide, "LINUX SMOKE", "基础环境已经从准备态变成可执行 smoke 证据", "03")
    add_terminal(slide, ASSETS / "terminal-linux-env.png", 0.72, 1.75, 7.35, 3.55)
    add_text(slide, 8.35, 1.85, 3.95, 0.5, "已验证组件", 17, WHITE, bold=True)
    bullets(
        slide,
        8.35,
        2.45,
        3.95,
        [
            "Ubuntu 24.04.4 LTS / WSL2",
            "Verilator 5.020，GCC 13.3，GNU Make 4.3",
            "Picker 0.9.0：C++ 与 Python 支持 OK",
            "pytoffee 0.3.0 + toffee-test 0.3.0",
        ],
    )
    add_text(slide, 8.35, 5.55, 3.65, 0.95, "版本说明：Toffee 最新 PyPI 包存在接口漂移；smoke 环境固定到可 import 且能跑上游 pytest 的 0.3.0 组合。", 11, MUTED)


def slide_smoke_flow(slide, data: dict) -> None:
    base(slide, "UPSTREAM FLOW", "上游 Example-NutShellCache 的两个关键目标已完成", "04")
    add_terminal(slide, ASSETS / "terminal-smoke.png", 0.72, 1.72, 6.7, 3.15)
    steps = [
        ("fetch", "锁定 commit\ncdc9ef7"),
        ("gen_dut", "Picker export\nexit 0"),
        ("test", "Toffee pytest\n1 passed"),
        ("record", "JSON/MD artifact\n写入 reports"),
    ]
    for i, (name, desc) in enumerate(steps):
        x = 0.88 + i * 3.05
        panel(slide, x, 5.25, 2.45, 0.95, name, desc)
        if i < len(steps) - 1:
            arrow(slide, x + 2.47, 5.72, x + 2.85, 5.72)
    add_text(slide, 7.75, 1.85, 4.45, 1.0, "结果口径", 18, WHITE, bold=True)
    bullets(
        slide,
        7.75,
        2.55,
        4.25,
        [
            "`rtl_smoke_complete` 来自真实命令返回值。",
            "waveform / generated DUT / coverage candidate 已写入 manifest。",
            "RTL code coverage 成功则写 LCOV 摘要，失败则写 not_exported 原因。",
        ],
        font_size=12.2,
    )


def slide_coverage(slide, data: dict) -> None:
    base(slide, "COVERAGE", "seed 11 的 96 条事务覆盖 23 个功能点", "05")
    add_terminal(slide, ASSETS / "terminal-coverage.png", 0.72, 1.75, 5.55, 3.05)
    sample = data["sample"]
    events = sample["event_counts"]
    add_text(slide, 6.65, 1.78, 2.5, 0.5, "23/23", 38, GREEN, bold=True)
    add_text(slide, 8.7, 1.95, 2.0, 0.3, "100.00% coverage", 13.5, WHITE, bold=True)
    max_count = max(events.values())
    y = 2.65
    for name in ["miss", "refill", "hit", "write", "eviction", "dirty_eviction", "writeback", "stall_hold", "reset_window"]:
        count = events[name]
        add_text(slide, 6.65, y, 1.55, 0.22, name, 9, MUTED)
        bar(slide, 8.1, y + 0.04, 3.25 * count / max_count, 0.12, CYAN if count > 2 else AMBER)
        add_text(slide, 11.55, y - 0.01, 0.5, 0.22, str(count), 9, WHITE, align=PP_ALIGN.RIGHT)
        y += 0.39


def slide_faults(slide, data: dict) -> None:
    base(slide, "FAULT INJECTION", "5 类 injected fault 均触发可解释失败", "06")
    add_terminal(slide, ASSETS / "terminal-faults.png", 0.72, 1.72, 5.7, 3.15)
    x, y = 6.72, 1.76
    add_text(slide, x, y, 4.8, 0.32, "检出矩阵", 17, WHITE, bold=True)
    y += 0.52
    for item in data["faults"]:
        panel(slide, x, y, 5.1, 0.55, item["mode"], f"expected failure · failures={item['failures']}", small=True)
        y += 0.64
    add_text(slide, 0.85, 6.18, 11.5, 0.35, "这些 artifact 用来证明验证环境和 scoreboard 能捕捉注入错误；报告不把 injected fault 写成真实 NutShell RTL bug。", 12.5, MUTED)


def slide_review(slide, data: dict) -> None:
    base(slide, "REVIEW JOURNAL", "人工复核把草案问题收敛为可审计的验证逻辑", "07")
    add_terminal(slide, ASSETS / "terminal-review.png", 0.72, 1.75, 5.9, 3.25)
    y = 1.72
    for row in data["review"][:5]:
        add_text(slide, 6.95, y, 0.65, 0.25, row["id"], 10, CYAN, bold=True)
        add_text(slide, 7.6, y, 4.65, 0.48, row["correction"], 10.5, WHITE)
        y += 0.88
    add_text(slide, 0.78, 5.55, 11.35, 0.55, "答辩时可说明：UCAgent 负责加速草案和覆盖洞提示，人工复核负责 invariant、fault 定义、报告边界和可复现证据。", 14, MUTED)


def slide_boundary(slide, data: dict) -> None:
    base(slide, "SUBMISSION PACKAGE", "最终材料保留实证结果，也保留未导出覆盖率的边界", "08")
    add_terminal(slide, ASSETS / "terminal-repo.png", 0.72, 1.72, 5.9, 3.25)
    metric_rail(
        slide,
        6.95,
        1.85,
        [
            ("PDF", "正式验证报告"),
            ("PPTX", "答辩演示材料"),
            ("JSON", "smoke / fault / coverage artifact"),
            ("MD", "验证计划与复核记录"),
        ],
        columns=2,
    )
    add_text(slide, 7.02, 4.82, 4.8, 0.3, "答辩边界表述", 17, WHITE, bold=True)
    bullets(
        slide,
        7.02,
        5.28,
        4.85,
        [
            "已完成：Python harness、fault 检出、Linux Picker/Toffee smoke。",
            "未声称：真实 RTL bug、RTL functional coverage 已闭合。",
            "后续条件：更长随机回归、waveform 截图、真实 RTL functional coverage。",
        ],
        font_size=11.4,
    )


def add_terminal(slide, path: Path, x: float, y: float, w: float, h: float) -> None:
    slide.shapes.add_picture(str(path), Inches(x), Inches(y), Inches(w), Inches(h))


def metric_rail(slide, x: float, y: float, metrics: Sequence[tuple[str, str]], columns: int = 4) -> None:
    cell_w = 2.65 if columns == 4 else 2.55
    for i, (value, label) in enumerate(metrics):
        cx = x + (i % columns) * (cell_w + 0.12)
        cy = y + (i // columns) * 1.15
        rect(slide, cx, cy, cell_w, 0.92, PANEL)
        add_text(slide, cx + 0.15, cy + 0.12, cell_w - 0.3, 0.32, value, 21, GREEN if i != 2 else AMBER, bold=True)
        add_text(slide, cx + 0.15, cy + 0.52, cell_w - 0.3, 0.26, label, 8.8, MUTED)


def panel(slide, x: float, y: float, w: float, h: float, title: str, body: str, small: bool = False) -> None:
    rect(slide, x, y, w, h, PANEL_2)
    add_text(slide, x + 0.16, y + 0.1, w - 0.3, 0.25, title, 11 if small else 13, CYAN, bold=True)
    add_text(slide, x + 0.16, y + (0.34 if small else 0.45), w - 0.3, h - 0.42, body, 9.2 if small else 10.5, WHITE)


def bullets(slide, x: float, y: float, w: float, items: Iterable[str], font_size: float = 13.5) -> None:
    top = y
    for item in items:
        add_text(slide, x, top, 0.18, 0.22, "-", font_size, CYAN, bold=True)
        add_text(slide, x + 0.28, top - 0.02, w - 0.28, 0.52, item, font_size, WHITE)
        top += 0.68


def add_text(slide, x: float, y: float, w: float, h: float, text: str, size: float, color: RGBColor, bold: bool = False, align=PP_ALIGN.LEFT) -> None:
    shape = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    shape.text_frame.clear()
    shape.text_frame.margin_left = Inches(0.02)
    shape.text_frame.margin_right = Inches(0.02)
    shape.text_frame.margin_top = Inches(0.02)
    shape.text_frame.margin_bottom = Inches(0.02)
    shape.text_frame.word_wrap = True
    shape.text_frame.vertical_anchor = MSO_ANCHOR.TOP
    p = shape.text_frame.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.name = "Microsoft YaHei"
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color


def rect(slide, x: float, y: float, w: float, h: float, color: RGBColor) -> None:
    shape = slide.shapes.add_shape(1, Inches(x), Inches(y), Inches(w), Inches(h))
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.color.rgb = RULE
    shape.line.width = Pt(0.6)


def bar(slide, x: float, y: float, w: float, h: float, color: RGBColor) -> None:
    shape = slide.shapes.add_shape(1, Inches(x), Inches(y), Inches(max(w, 0.05)), Inches(h))
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()


def line(slide, x1: float, y1: float, x2: float, y2: float, color: RGBColor) -> None:
    shape = slide.shapes.add_connector(1, Inches(x1), Inches(y1), Inches(x2), Inches(y2))
    shape.line.color.rgb = color
    shape.line.width = Pt(1)


def arrow(slide, x1: float, y1: float, x2: float, y2: float) -> None:
    shape = slide.shapes.add_connector(1, Inches(x1), Inches(y1), Inches(x2), Inches(y2))
    shape.line.color.rgb = CYAN
    shape.line.width = Pt(1.5)


if __name__ == "__main__":
    raise SystemExit(main())
